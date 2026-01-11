import fitz 
from pptx import Presentation
import os

def load_pdf(file_path):
    """Extracts text from a PDF file."""
    try:
        doc = fitz.open(file_path)
        text = []
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            if page_text.strip():
                text.append(f"--- Page {page_num + 1} ---\n{page_text}")
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return None

def load_pptx(file_path):
    """Extracts text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return None


def load_document(file_path):
    """Wrapper to detect file extension and call appropriate loader."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return load_pdf(file_path)
    elif ext == '.pptx':
        return load_pptx(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return None

# Test the functions
if __name__ == "__main__":
    print("Ingestion module loaded.")
    # file_path = "data/notes.docx"
    # print(load_document(file_path))